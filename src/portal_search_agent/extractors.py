from __future__ import annotations

import csv
import html
import io
import re
from pathlib import Path
from typing import Iterable
from urllib.parse import urlparse

import httpx
from bs4 import BeautifulSoup
from charset_normalizer import from_bytes
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from striprtf.striprtf import rtf_to_text

from .config import Settings
from .models import LinkCandidate
from .urltools import is_probably_document, normalize_url, should_skip_url


class ExtractionError(RuntimeError):
    pass


def html_to_text(html: bytes | str) -> tuple[str, str]:
    if isinstance(html, bytes):
        html = decode_bytes(html)
    soup = BeautifulSoup(html, "lxml")
    for node in soup(["script", "style", "noscript", "template"]):
        node.decompose()

    title = ""
    if soup.title and soup.title.string:
        title = clean_text(soup.title.string)
    if not title:
        h1 = soup.find("h1")
        if h1:
            title = clean_text(h1.get_text(" ", strip=True))

    text = clean_text(soup.get_text(" ", strip=True))
    return title, text


def extract_links_from_html(html: bytes | str, base_url: str) -> list[str]:
    if isinstance(html, bytes):
        html = decode_bytes(html)
    soup = BeautifulSoup(html, "lxml")
    links: list[str] = []
    for selector, attr in (("a[href]", "href"), ("area[href]", "href"), ("iframe[src]", "src"), ("frame[src]", "src")):
        for node in soup.select(selector):
            value = node.get(attr)
            normalized = normalize_url(value or "", base_url)
            if normalized:
                links.append(normalized)
    return list(dict.fromkeys(links))


def extract_actionable_links_from_html(html: bytes | str, base_url: str, source_title: str = "") -> list[LinkCandidate]:
    if isinstance(html, bytes):
        html = decode_bytes(html)
    soup = BeautifulSoup(html, "lxml")
    links: list[LinkCandidate] = []
    seen: set[tuple[str, str]] = set()

    for selector, attr in (("a[href]", "href"), ("area[href]", "href"), ("iframe[src]", "src"), ("frame[src]", "src")):
        for node in soup.select(selector):
            raw_url = node.get(attr) or ""
            normalized = normalize_url(raw_url, base_url, keep_fragment=True)
            if not normalized or should_skip_url(normalized):
                continue

            label = link_label(node, normalized)
            context = link_context(node, label, source_title)
            if len((label + context).strip()) < 2:
                continue

            key = (normalized, label.lower())
            if key in seen:
                continue
            seen.add(key)

            links.append(
                LinkCandidate(
                    url=normalized,
                    text=label,
                    context=context,
                    kind=classify_link(node, normalized, label, context),
                    source_url=base_url,
                    source_title=source_title,
                )
            )

    return links


def link_label(node, url: str) -> str:
    parts = [
        node.get_text(" ", strip=True),
        node.get("aria-label", ""),
        node.get("title", ""),
    ]
    for image in node.find_all("img"):
        parts.append(image.get("alt", ""))
        parts.append(image.get("title", ""))

    label = clean_text(" ".join(part for part in parts if part))
    if label:
        return label[:240]

    path_name = Path(urlparse(url).path).name
    if path_name:
        return clean_text(path_name.replace("-", " ").replace("_", " "))[:240]
    return url


def link_context(node, label: str, source_title: str = "") -> str:
    interesting_tags = {"nav", "li", "tr", "td", "th", "section", "article", "header", "footer"}
    interesting_class_parts = ("nav", "menu", "tab", "app", "application", "card", "tile", "portal")

    for parent in node.parents:
        if getattr(parent, "name", None) in {"body", "html", "[document]"}:
            break

        class_text = " ".join(parent.get("class", [])).lower() if hasattr(parent, "get") else ""
        role = (parent.get("role", "") if hasattr(parent, "get") else "").lower()
        is_interesting = (
            getattr(parent, "name", None) in interesting_tags
            or role in {"navigation", "menu", "menubar", "tablist", "tab"}
            or any(part in class_text for part in interesting_class_parts)
        )
        if not is_interesting:
            continue

        text = clean_text(parent.get_text(" ", strip=True))
        if text and text.lower() != label.lower():
            return text[:500]

    return source_title[:500] if source_title else ""


def classify_link(node, url: str, label: str, context: str) -> str:
    role = (node.get("role", "") if hasattr(node, "get") else "").lower()
    class_text = " ".join(node.get("class", [])).lower() if hasattr(node, "get") else ""
    combined = f"{label} {context} {url}".lower()

    if urlparse(url).fragment:
        return "tab_link"
    if is_probably_document(url):
        return "document_link"
    if any(marker in combined for marker in ("app", "application", "εφαρμογ")):
        return "application_link"
    if role == "tab" or "tab" in class_text or "#tab" in url.lower():
        return "tab_link"
    if role in {"menuitem", "navigation"} or any(marker in class_text for marker in ("nav", "menu")):
        return "navigation_link"
    return "portal_link"


def extract_sitemap_urls(xml: bytes | str) -> list[str]:
    if isinstance(xml, bytes):
        xml = decode_bytes(xml)
    return [clean_text(html.unescape(match)) for match in re.findall(r"<loc>\s*(.*?)\s*</loc>", xml, flags=re.I | re.S)]


async def extract_with_tika(data: bytes, content_type: str, filename: str, settings: Settings) -> str:
    suffix = Path(filename).suffix.lower()
    headers = {
        "Accept": "text/plain",
        "Content-Type": content_type or "application/octet-stream",
        "X-Tika-Timeout-Millis": "300000",
    }
    if content_type == "application/pdf" or suffix == ".pdf" or content_type.startswith("image/"):
        headers["X-Tika-OCRLanguage"] = settings.ocr_languages
        headers["X-Tika-PDFextractInlineImages"] = "true"
    else:
        headers["X-Tika-OCRskipOcr"] = "true"

    url = f"{settings.tika_url}/tika"
    async with httpx.AsyncClient(timeout=settings.request_timeout_seconds * 5) as client:
        response = await client.put(url, content=data, headers=headers)
    if response.status_code >= 400:
        raise ExtractionError(f"Tika extraction failed for {filename}: HTTP {response.status_code}")
    return clean_text(response.text)


async def extract_document(data: bytes, content_type: str, filename: str, settings: Settings) -> str:
    try:
        text = await extract_with_tika(data, content_type, filename, settings)
        if text:
            return text
    except Exception:
        pass

    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        return extract_pdf_fallback(data)
    if suffix == ".docx":
        return extract_docx_fallback(data)
    if suffix == ".xlsx":
        return extract_xlsx_fallback(data)
    if suffix == ".pptx":
        return extract_pptx_fallback(data)
    if suffix == ".rtf":
        return clean_text(rtf_to_text(decode_bytes(data)))
    if suffix in {".txt", ".csv", ".xml"} or content_type.startswith("text/"):
        return extract_text_fallback(data, suffix)
    return clean_text(decode_bytes(data))


def extract_pdf_fallback(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    parts = [page.extract_text() or "" for page in reader.pages]
    return clean_text("\n".join(parts))


def extract_docx_fallback(data: bytes) -> str:
    document = DocxDocument(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append(" ".join(cell.text for cell in row.cells))
    return clean_text("\n".join(parts))


def extract_xlsx_fallback(data: bytes) -> str:
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(sheet.title)
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                parts.append(" ".join(values))
    workbook.close()
    return clean_text("\n".join(parts))


def extract_pptx_fallback(data: bytes) -> str:
    deck = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in deck.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return clean_text("\n".join(parts))


def extract_text_fallback(data: bytes, suffix: str) -> str:
    text = decode_bytes(data)
    if suffix == ".csv":
        rows = csv.reader(io.StringIO(text))
        return clean_text("\n".join(" ".join(row) for row in rows))
    return clean_text(text)


def decode_bytes(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        pass

    match = from_bytes(data).best()
    if match:
        return str(match)
    return data.decode("utf-8", errors="ignore")


def clean_text(value: str | Iterable[str]) -> str:
    if not isinstance(value, str):
        value = " ".join(str(item) for item in value)
    return re.sub(r"\s+", " ", value).strip()
